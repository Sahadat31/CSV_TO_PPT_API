import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from textwrap import wrap

def clean_markdown(text: str) -> str:
    return re.sub(r'(\*{1,2})(.*?)\1', r'\2', text).replace("`", "").strip()

def add_wrapped_text_slide(prs, title: str, content: str, max_chars=800):
    """Split long content into multiple slides if it exceeds character limit."""
    chunks = wrap(content.strip(), width=max_chars, break_long_words=False)
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title = f"{title} (Part {i+1})" if len(chunks) > 1 else title
        slide.shapes.title.text = slide_title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = clean_markdown(chunk)
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT

def add_single_text_slide(prs, title: str, content: str):
    """One text slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = clean_markdown(content)
    p.font.size = Pt(18)
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

def add_chart_slide(prs, title: str, image_path: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(4), Inches(1), width=Inches(8))

def add_thank_you_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Thank  You'

def add_bullet_slide(prs, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for bullet in bullets:
        if bullet.strip():
            p = tf.add_paragraph()
            p.text = clean_markdown(bullet.strip("-• ").strip())
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = "Calibri"

def create_ppt(summary: str, insights: list[str], chart_paths: list[str]) -> str:
    prs = Presentation("assets/template.pptx")

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI-Generated Analysis Report"
    slide.placeholders[1].text = "Structured from unstructured CSV data"

    # Summary Slides (break into multiple if too long)
    add_wrapped_text_slide(prs, "Executive Summary", summary)

    # Insight Slides (each item → 1 slide)
    for i, insight in enumerate(insights):
        title = f"Insight {i+1}"
        insight.replace('*','')
        add_wrapped_text_slide(prs, title, insight)

    # Charts Section
    for i, path in enumerate(chart_paths):
        add_chart_slide(prs, f"Chart {i+1}", path)

    # Closing Recommendations (optional static content)
    recommendations = [
        "Ensure broader geographic sampling in future surveys.",
        "Review screening logic (S1–S4) for potential over-filtering.",
        "Explore battery performance differences between brands.",
        "Investigate time inefficiencies in patient transport.",
        "Continue improving comfort and safety feedback collection."
    ]
    add_bullet_slide(prs, "AI Recommendations", recommendations)
    add_thank_you_slide(prs)

    output_path = "output/generated_presentation_structured.pptx"
    prs.save(output_path)
    return output_path
